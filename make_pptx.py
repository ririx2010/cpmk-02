#!/usr/bin/env python3
"""
Generate PPTX presentation for Ujian CPMK-2.
Each slide has speaker notes as presentation hints.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTDIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD   = RGBColor(0x27, 0x27, 0x44)
ACCENT    = RGBColor(0xE9, 0x45, 0x60)
ACCENT2   = RGBColor(0x0F, 0x34, 0x60)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
GREEN     = RGBColor(0x00, 0xC8, 0x53)
BLUE      = RGBColor(0x53, 0x8D, 0xED)
RED       = RGBColor(0xE9, 0x45, 0x60)

def add_bg(slide, color=BG_DARK):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT, spacing=Pt(8)):
    """Add bulleted list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_notes(slide, text):
    """Add speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_card(slide, left, top, width, height, color=BG_CARD):
    """Add rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width):
    """Add accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

# ═════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_accent_line(slide, 1, 2.8, 4)

add_textbox(slide, 1, 1.2, 11, 1.5,
    "Sistem Pendukung Keputusan\nPemilihan Bidang Konsentrasi",
    font_size=36, bold=True, color=WHITE)
add_textbox(slide, 1, 3.0, 11, 0.8,
    "Fuzzy Inference System — Metode Mamdani",
    font_size=22, color=GOLD)
add_textbox(slide, 1, 4.5, 11, 1.5,
    "RIYADI — H2A025002\nMagister Teknik Elektro, Universitas Jenderal Soedirman\nUjian CPMK-2 | Kecerdasan Buatan (MTE25112)",
    font_size=16, color=LIGHT)

add_notes(slide, """HINT — SLIDE 1 (Judul) — ~15 detik

Assalamualaikum / Selamat pagi Bapak/Ibu dosen.
Perkenalkan, saya Riyadi, NIM H2A025002.
Pada kesempatan ini saya akan mempresentasikan Ujian CPMK-2:
rancangan Sistem Pendukung Keputusan untuk pemilihan Bidang Konsentrasi
menggunakan Fuzzy Inference System metode Mamdani.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 2: Masalah
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Latar Belakang Masalah", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

add_bullet_list(slide, 0.8, 1.4, 6, 5.5, [
    "Perguruan tinggi memiliki 3 KBK: ITP, RPLD, SKJK",
    "Setiap mahasiswa mengisi survei: IPK dan Minat (skala 0-4)",
    "Sistem harus menghitung Nilai Kelayakan (NK) per KBK",
    "KBK dengan NK tertinggi = rekomendasi",
    "",
    "DATA TARGET #20:",
    "  ITP:  IP=3.40, Minat=2.60",
    "  RPLD: IP=3.52, Minat=2.80",
    "  SKJK: IP=3.65, Minat=1.50",
    "",
    "Pertanyaan: KBK mana yang direkomendasikan?",
    "  → Syarat: HARUS RPLD"
], font_size=17, color=LIGHT)

# Add decision boundary plot on right
plot_path = os.path.join(OUTDIR, "plot_decision_boundary-1.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(7.2), Inches(1.2), Inches(5.5), Inches(4.2))

add_card(slide, 7.2, 5.6, 5.5, 1.3, BG_CARD)
add_textbox(slide, 7.4, 5.7, 5.1, 1.1,
    "Heatmap: IP tinggi saja tidak cukup!\nSKJK punya IP tertinggi (3.65) tapi Minat terendah (1.50)\n→ Fuzzy bisa menangkap kedua faktor secara bersamaan",
    font_size=13, color=GOLD)

add_notes(slide, """HINT — SLIDE 2 (Masalah) — ~1 menit

Masalahnya: sebuah kampus ingin membangun SPK untuk membantu mahasiswa
memilih konsentrasi di antara 3 pilihan: ITP, RPLD, dan SKJK.

Setiap mahasiswa punya data IPK dan Minat untuk setiap KBK.
Sistem fuzzy menghitung Nilai Kelayakan (NK) untuk masing-masing KBK,
lalu merekomendasikan yang tertinggi.

Data target nomor 20 menarik karena:
- SKJK punya IP TERTINGGI (3.65) tapi minat paling rendah (1.50)
- RPLD punya kombinasi paling seimbang
- Pertanyaannya: apakah fuzzy system bisa menangkap nuansa ini?

Lihat heatmap di kanan: area merah = NK tinggi, butuh IP DAN Minat
yang tinggi secara bersamaan. Ini menunjukkan kenapa kita butuh fuzzy,
bukan sekadar mengambil IP tertinggi.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 3: Arsitektur Mamdani
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Arsitektur FIS Mamdani", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# 5 step boxes
steps = [
    ("1", "Fuzzifikasi", "Crisp → μ\nIP=3.52 → μ=0.96", BLUE),
    ("2", "Evaluasi", "α = min(μ₁, μ₂)\nα = 0.80", RGBColor(0x4E, 0x9A, 0xD1)),
    ("3", "Implikasi", "MF dipotong\npada α=0.80", RGBColor(0xE8, 0xA8, 0x38)),
    ("4", "Agregasi", "max semua\ntrapesium", RGBColor(0xE9, 0x70, 0x50)),
    ("5", "Defuzzifikasi", "COG =\ntitik pusat", GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.6 + i * 2.5
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.5), Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Card
    card = add_card(slide, x, 2.3, 2.2, 2.2, BG_CARD)
    add_textbox(slide, x + 0.1, 2.4, 2.0, 0.5, title, font_size=18, bold=True, color=color)
    add_textbox(slide, x + 0.1, 2.9, 2.0, 1.5, desc, font_size=14, color=LIGHT)

    # Arrow (except last)
    if i < 4:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(x + 2.2), Inches(3.0), Inches(0.35), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Bottom example
add_card(slide, 0.6, 5.0, 12.0, 2.0, BG_CARD)
add_textbox(slide, 0.8, 5.1, 11.5, 0.5, "Contoh: Data #20 → KBK RPLD (IP=3.52, Minat=2.80)",
    font_size=16, bold=True, color=GOLD)
add_textbox(slide, 0.8, 5.6, 11.5, 1.2,
    "Input crisp (3.52, 2.80) → Fuzzifikasi: μ_Bagus=0.96, μ_Suka=0.80 → "
    "Evaluasi: α=min(0.96, 0.80)=0.80 → Implikasi: MF Tinggi dipotong di 0.80 → "
    "Agregasi: 1 trapesium → Defuzzifikasi COG: NK = 75.00",
    font_size=14, color=LIGHT)

add_notes(slide, """HINT — SLIDE 3 (Arsitektur Mamdani) — ~45 detik

Ini diagram alir FIS Mamdani yang saya gunakan. Ada 5 tahap:

1. FUZZIFIKASI: mengubah angka crisp menjadi derajat keanggotaan.
   Misal IP 3.52 → μ_Bagus = 0.96 (hampir penuh)

2. EVALUASI ATURAN: setiap rule dihitung kekuatan aktivasi α = min(μ_IP, μ_Minat).
   Jadi α = min(0.96, 0.80) = 0.80

3. IMPLIKASI: MF output (Tinggi) dipotong pada tinggi α=0.80, jadi trapesium.

4. AGREGASI: semua trapesium digabung dengan operator max.

5. DEFUZZIFIKASI: hitung Centre of Gravity (COG) = titik pusat massa dari area.
   Untuk RPLD data #20, COG = 75.00

Contoh di bawah menunjukkan end-to-end untuk data RPLD #20.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 4: Fuzzifikasi
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Fuzzifikasi — Membership Functions", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# Add MF detail plot
plot_path = os.path.join(OUTDIR, "plot_mf_detail-1.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5))

# Bottom cards
add_card(slide, 0.5, 6.0, 4.0, 1.2, BG_CARD)
add_textbox(slide, 0.7, 6.05, 3.6, 1.1,
    "IP: Kecil(0,1,2) Cukup(1.5,2.5,3.5) Bagus(3,3.5,4)\n"
    "Kenapa segitiga? Sederhana, efisien, cukup\nuntuk merepresentasikan gradasi akademik.",
    font_size=12, color=LIGHT)

add_card(slide, 4.7, 6.0, 4.0, 1.2, BG_CARD)
add_textbox(slide, 4.9, 6.05, 3.6, 1.1,
    "Minat: TS(0,1.5,2.5) Suka(2,3,4) SS(3,3.5,4)\n"
    "Suka berpuncak di 3, bukan 2.5 — sehingga\nminat 2.80 > 2.60 dalam derajat μ",
    font_size=12, color=LIGHT)

add_card(slide, 8.9, 6.0, 3.9, 1.2, BG_CARD)
add_textbox(slide, 9.1, 6.05, 3.5, 1.1,
    "Segitiga dipilih karena:\n• Komputasi ringan\n• Overlap halus antar himpunan\n"
    "• Cocok untuk variabel kontinu",
    font_size=12, color=GOLD)

add_notes(slide, """HINT — SLIDE 4 (Fuzzifikasi) — ~1.5 menit

Perhatikan grafik di atas. Dua variabel input: IP (kiri) dan Minat (kanan).

Titik-titik warna menunjukkan posisi data #20 untuk tiap KBK:
- Biru = ITP, Merah = RPLD, Hijau = SKJK

PENTING — lihat grafik Minat (kanan):
- SKJK (hijau) jatuh tepat di puncak "Tidak Suka" → μ = 1.0
- ITP (biru) di μ_Suka = 0.6
- RPLD (merah) di μ_Suka = 0.8 → TERTINGGI

Ini sudah menunjukkan kenapa RPLD menang: minatnya paling tinggi
di antara ketiganya pada himpunan "Suka".

Kenapa Suka berpuncak di 3? Karena kalau puncaknya di 2.5,
maka minat 2.60 dan 2.80 akan hampir sama μ-nya.
Dengan puncak di 3, selisih 0.6 vs 0.8 lebih signifikan.

Kenapa segitiga, bukan trapesium? Karena untuk variabel akademik,
representasi linear sudah memadai dan komputasinya lebih ringan.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 5: 9 Fuzzy Rules
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "9 Fuzzy Rules", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# Rule matrix as cards
rules_matrix = [
    # Row labels + 3 columns
    ["Kecil",  "Rendah\n(R1)", "Rendah\n(R2)", "Sedang\n(R3)"],
    ["Cukup",  "Rendah\n(R4)", "Sedang\n(R5)", "Tinggi\n(R6)"],
    ["Bagus",  "Sedang\n(R7)", "Tinggi\n(R8)", "Tinggi\n(R9)"],
]

nk_colors = {"Rendah": RGBColor(0xE9, 0x45, 0x60),
             "Sedang": RGBColor(0xE8, 0xA8, 0x38),
             "Tinggi": RGBColor(0x00, 0xC8, 0x53)}

# Headers
for j, h in enumerate(["IP \\ Minat", "Tidak Suka", "Suka", "Sangat Suka"]):
    x = 1.5 + j * 2.7
    add_card(slide, x, 1.4, 2.4, 0.6, ACCENT2 if j > 0 else BG_CARD)
    add_textbox(slide, x + 0.1, 1.45, 2.2, 0.5, h, font_size=14, bold=True,
                color=WHITE if j > 0 else LIGHT, alignment=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(rules_matrix):
    for j, cell in enumerate(row):
        x = 1.5 + j * 2.7
        y = 2.2 + i * 1.2
        bg_color = BG_CARD
        txt_color = LIGHT
        nk_key = None

        if j == 0:
            txt_color = WHITE
        else:
            for k in nk_colors:
                if k in cell:
                    nk_key = k
                    txt_color = nk_colors[k]
                    break

        # Highlight R7
        if i == 2 and j == 1:
            bg_color = RGBColor(0x4A, 0x1A, 0x2E)

        # Highlight R8
        if i == 2 and j == 2:
            bg_color = RGBColor(0x0A, 0x3A, 0x2E)

        add_card(slide, x, y, 2.4, 0.9, bg_color)
        add_textbox(slide, x + 0.1, y + 0.05, 2.2, 0.8, cell,
                    font_size=15, bold=(j == 0), color=txt_color,
                    alignment=PP_ALIGN.CENTER)

# Right side: key insight
add_card(slide, 8.5, 1.4, 4.3, 5.5, BG_CARD)
add_textbox(slide, 8.7, 1.5, 3.9, 0.5, "Prinsip Desain", font_size=18, bold=True, color=GOLD)
add_bullet_list(slide, 8.7, 2.1, 3.9, 4.5, [
    "IP rendah → max Sedang\n(syarat akademik belum terpenuhi)",
    "",
    "Minat rendah → batasi output\n(ketidaksesuaian minat = risiko)",
    "",
    "R7 = DIFFERENTIATOR KUNCI:\nBagus × Tidak Suka → Sedang\n(bukan Tinggi!)",
    "",
    "R8 = RULE UTAMA:\nBagus × Suka → Tinggi\n(ini yang membuat RPLD menang)",
], font_size=13, color=LIGHT)

add_notes(slide, """HINT — SLIDE 5 (9 Fuzzy Rules) — ~1.5 menit

Ini matriks 3×3 dari 9 fuzzy rules. Baris = IP, Kolom = Minat.

Warna merah = Rendah, kuning = Sedang, hijau = Tinggi.

Perhatikan 2 rule yang paling penting:

RULE R7 (highlight ungu) — Bagus × Tidak Suka → Sedang, BUKAN Tinggi.
Kenapa? Karena pengalaman menunjukkan: nilai bagus tapi tidak berminat
sering berujung drop-out atau ketidakpuasan. Jadi saya batasi hanya Sedang.
INI ADALAH KUNCI yang membuat SKJK kalah di data #20:
meskipun IP tertinggi (3.65), minatnya hanya 1.50 (Tidak Suka),
sehingga R7 membatasi NK-nya jadi Sedang (50).

RULE R8 (highlight hijau) — Bagus × Suka → Tinggi.
Ini rule yang menguntungkan RPLD: IP 3.52 (Bagus) dan Minat 2.80 (Suka)
menghasilkan α = 0.80, output Tinggi, NK = 75.

Prinsip desainnya: baik IP MAUPUN Minat harus memadai untuk mendapat
rekomendasi Tinggi. Tidak cukup hanya salah satu.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 6: Kenapa Mamdani?
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Pemilihan Metode Inferensi", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# 3 method comparison cards
methods = [
    ("Mamdani", "DIPILIH", GREEN, [
        "Output: area fuzzy (trapesium)",
        "Defuzzifikasi: COG (integral)",
        "BISA divisualisasikan",
        "Cocok untuk SPK (transparan)",
        "Konsisten dengan materi kuliah",
    ]),
    ("Sugeno", "Alternatif", LIGHT, [
        "Output: fungsi linear/konstanta",
        "Defuzzifikasi: weighted average",
        "Lebih simpel, tapi TANPA visual",
        "Cocok untuk kontrol real-time",
        "Kurang intuitif untuk SPK",
    ]),
    ("Tsukamoto", "Alternatif", LIGHT, [
        "Output: titik per rule (invers MF)",
        "Defuzzifikasi: weighted average",
        "Perlu MF monotonik",
        "Lebih simpel dari Mamdani",
        "Tapi kurang fleksibel",
    ]),
]

for i, (name, badge, badge_color, items) in enumerate(methods):
    x = 0.5 + i * 4.2
    bg = BG_CARD if i > 0 else RGBColor(0x0A, 0x3A, 0x2E)
    add_card(slide, x, 1.4, 3.9, 5.5, bg)
    add_textbox(slide, x + 0.2, 1.5, 3.5, 0.5, name, font_size=22, bold=True, color=WHITE)
    add_textbox(slide, x + 0.2, 2.1, 3.5, 0.4, badge, font_size=14, bold=True, color=badge_color)
    add_bullet_list(slide, x + 0.2, 2.7, 3.5, 3.5, items, font_size=13, color=LIGHT, spacing=Pt(10))

add_notes(slide, """HINT — SLIDE 6 (Pemilihan Metode) — ~45 detik

Saya membandingkan 3 metode: Mamdani, Sugeno, dan Tsukamoto.

Saya memilih MAMDANI karena 3 alasan:

1. BISA DIVISUALISASIKAN — area fuzzy output (trapesium) bisa digambar.
   Sugeno dan Tsukamoto hanya menghasilkan angka, tidak ada visual.
   Untuk SPK, transparansi sangat penting.

2. COCOK UNTUK SPK — pengambil keputusan perlu memahami MENGAPA
   suatu rekomendasi diberikan. Dengan Mamdani, kita bisa trace
   setiap langkah dari fuzzifikasi hingga COG.

3. KONSISTEN — metode ini sama dengan yang dipelajari di materi
   perkuliahan dan tugas mandiri pertemuan 7.

Kalau ditanya "kenapa tidak Sugeno?":
Sugeno memang lebih simpel secara komputasi, tapi tidak menghasilkan
area fuzzy. Untuk SPK yang perlu dijelaskan ke stakeholder,
Mamdani lebih tepat karena bisa divisualisasikan.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 7: Walkthrough Data #20 (THE KEY SLIDE)
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.3, 11, 0.7, "Walkthrough: Data Target #20", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 0.95, 3)

# 3 KBK cards side by side
kbk_data = [
    ("ITP", "IP=3.40  Minat=2.60", BLUE,
     "μ_Cukup=0.1, μ_Bagus=0.8\nμ_Suka=0.6",
     "R5: α=min(0.1, 0.6) = 0.1\nR8: α=min(0.8, 0.6) = 0.6",
     "NK = 71.14", "Sedang@0.1 + Tinggi@0.6"),
    ("RPLD", "IP=3.52  Minat=2.80", RED,
     "μ_Bagus=0.96\nμ_Suka=0.80",
     "R8: α=min(0.96, 0.80) = 0.80",
     "NK = 75.00 ★", "Tinggi@0.8 (murni)"),
    ("SKJK", "IP=3.65  Minat=1.50", RGBColor(0x50, 0xC8, 0x78),
     "μ_Bagus=0.7\nμ_TS=1.0 (puncak!)",
     "R7: α=min(0.7, 1.0) = 0.70",
     "NK = 50.00", "Sedang@0.7 (terbatas)"),
]

for i, (kbk, input_str, color, fuzz, rules, nk, note) in enumerate(kbk_data):
    x = 0.4 + i * 4.2
    add_card(slide, x, 1.2, 4.0, 5.8, BG_CARD)

    # KBK name
    add_textbox(slide, x + 0.2, 1.3, 3.6, 0.5, kbk, font_size=24, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 1.8, 3.6, 0.4, input_str, font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Fuzzifikasi
    add_textbox(slide, x + 0.2, 2.3, 3.6, 0.3, "Fuzzifikasi:", font_size=12, bold=True, color=GOLD)
    add_textbox(slide, x + 0.2, 2.6, 3.6, 0.8, fuzz, font_size=12, color=LIGHT)

    # Rules
    add_textbox(slide, x + 0.2, 3.5, 3.6, 0.3, "Rules Aktif:", font_size=12, bold=True, color=GOLD)
    add_textbox(slide, x + 0.2, 3.8, 3.6, 0.9, rules, font_size=12, color=LIGHT)

    # Note
    add_textbox(slide, x + 0.2, 4.8, 3.6, 0.3, "Implikasi:", font_size=12, bold=True, color=GOLD)
    add_textbox(slide, x + 0.2, 5.1, 3.6, 0.5, note, font_size=12, color=LIGHT)

    # NK result
    nk_color = GREEN if "★" in nk else LIGHT
    add_textbox(slide, x + 0.2, 5.8, 3.6, 0.6, nk, font_size=22, bold=True, color=nk_color, alignment=PP_ALIGN.CENTER)

# Bottom conclusion
add_card(slide, 0.4, 7.1, 12.5, 0.3, RGBColor(0x0A, 0x3A, 0x2E))

add_notes(slide, """HINT — SLIDE 7 (Walkthrough #20) — ~3 menit — SLIDE TERPENTING

Ini slide paling penting. Saya akan jelaskan step-by-step.

=== KBK ITP (kiri) ===
Input: IP=3.40, Minat=2.60
Fuzzifikasi: IP masuk Cukup (μ=0.1) dan Bagus (μ=0.8).
             Minat masuk Suka (μ=0.6).
Rules aktif: R5 (Cukup×Suka→Sedang, α=0.1) dan R8 (Bagus×Suka→Tinggi, α=0.6).
Hasil: Tinggi dominan, tapi ada sedikit Sedang yang menarik COG turun.
NK = 71.14

=== KBK RPLD (tengah) — PEMENANG ===
Input: IP=3.52, Minat=2.80
Fuzzifikasi: IP masuk Bagus (μ=0.96 — hampir penuh!).
             Minat masuk Suka (μ=0.80 — lebih tinggi dari ITP).
Hanya 1 rule aktif: R8 (Bagus×Suka→Tinggi, α=0.80).
Karena trapesium simetris, COG = titik puncak = 75.00

=== KBK SKJK (kanan) ===
Input: IP=3.65, Minat=1.50
Fuzzifikasi: IP masuk Bagus (μ=0.7). Minat masuk TS (μ=1.0 — puncak!).
Hanya 1 rule aktif: R7 (Bagus×TS→Sedang, α=0.70).
R7 adalah differentiator kunci: nilai bagus TANPA minat → hanya Sedang!
NK = 50.00

KESIMPULAN: RPLD (75) > ITP (71.14) > SKJK (50) → RPLD MENANG!

Point penjelasan: meskipun SKJK punya IP TERTINGGI (3.65),
rule R7 membatasinya karena minatnya sangat rendah.
Ini menunjukkan fuzzy system menangkap kedua faktor secara seimbang.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 8: Hasil Verifikasi
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Hasil Verifikasi — 20 Mahasiswa", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# Bar chart
plot_path = os.path.join(OUTDIR, "plot_bar_all-1.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(4.3))

# Bottom stats
add_card(slide, 0.5, 5.8, 3.8, 1.4, RGBColor(0x0A, 0x3A, 0x2E))
add_textbox(slide, 0.7, 5.85, 3.4, 1.3,
    "Akurasi Data Latih\n19/19 = 100%",
    font_size=20, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

add_card(slide, 4.6, 5.8, 4.0, 1.4, RGBColor(0x0A, 0x3A, 0x2E))
add_textbox(slide, 4.8, 5.85, 3.6, 1.3,
    "Data Target #20\nRPLD terpilih ✓",
    font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

add_card(slide, 8.9, 5.8, 3.9, 1.4, BG_CARD)
add_textbox(slide, 9.1, 5.85, 3.5, 1.3,
    "Verifikasi Python\nsim_fuzzy_spk.py\n10.001 titik diskritisasi",
    font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_notes(slide, """HINT — SLIDE 8 (Verifikasi) — ~1.5 menit

Grafik bar menunjukkan hasil NK untuk semua 20 mahasiswa.
3 warna: biru=ITP, merah=RPLD, hijau=SKJK.
Bintang menandakan KBK yang direkomendasikan.

Perhatikan:
- Data 1-7 (pemilih ITP): bar biru selalu tertinggi ✓
- Data 8-13 (pemilih RPLD): bar merah selalu tertinggi ✓
- Data 14-19 (pemilih SKJK): bar hijau selalu tertinggi ✓
- Data 20 (TARGET): bar merah (RPLD) tertinggi ✓

Hasil: akurasi 100% pada 19 data latih, dan data target #20
berhasil merekomendasikan RPLD sesuai syarat soal ujian.

Semua perhitungan diverifikasi menggunakan Python
dengan 10.001 titik diskritisasi untuk akurasi numerik.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 9: Decision Boundary (with explanation)
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Decision Boundary — Visualisasi Global", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

# Decision boundary plot (left side)
plot_path = os.path.join(OUTDIR, "plot_decision_boundary-1.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.3), Inches(1.3), Inches(7.5), Inches(5.7))

# Explanation card (right side)
add_card(slide, 8.0, 1.3, 4.8, 5.7, BG_CARD)
add_textbox(slide, 8.2, 1.4, 4.4, 0.5, "Apa itu Decision Boundary?",
            font_size=16, bold=True, color=GOLD)

add_bullet_list(slide, 8.2, 1.9, 4.4, 5.0, [
    "Peta warna yang menjawab:",
    "\"Jika IPK dan Minat sekian, NK-nya berapa?\"",
    "",
    "CARA MEMBACA:",
    "• Sumbu X = IPK, Sumbu Y = Minat",
    "• Merah = NK rendah (0-50)",
    "• Kuning = NK sedang (50-75)",
    "• Hijau = NK tinggi (75-100)",
    "",
    "DATA TARGET #20:",
    "○ RPLD → zona HIJAU (NK=75)",
    "  Kombinasi IP+Minat terbaik",
    "□ ITP → zona KUNING-HIJAU (NK=71)",
    "  Cukup baik, minat kurang tinggi",
    "△ SKJK → zona MERAH (NK=50)",
    "  IP tinggi tapi Minat sangat rendah",
    "",
    "INTI: NK tinggi HANYA tercapai jika",
    "IPK dan Minat KEDUANYA memadai.",
    "Area kanan-bawah tetap MERAH",
    "→ konsekuensi Rule R7!",
], font_size=11, color=LIGHT, spacing=Pt(3))

add_notes(slide, """HINT — SLIDE 9 (Decision Boundary) — ~1.5 menit

Ini adalah decision boundary chart yang sangat penting untuk memahami
secara GLOBAL bagaimana sistem fuzzy membuat keputusan.

BAGAIMANA MEMBACANYA:
- Ini seperti peta cuaca: warna menunjukkan NK untuk setiap kombinasi (IPK, Minat)
- Sumbu X = IPK, Sumbu Y = Minat
- Merah = NK rendah, kuning = NK sedang, hijau = NK tinggi
- Garis putus-putus = batas threshold NK=50 dan NK=75

INTERPRETASI DATA #20:
- RPLD (titik merah) berada di zona HIJAU → kombinasi terbaik
- ITP (titik biru) di zona KUNING-HIJAU → cukup tapi bukan tertinggi
- SKJK (titik hijau) di zona MERAH → meskipun IP TERTINGGI, minat terlalu rendah

POINT TERPENTING:
Perhatikan area KANAN BAWAH chart (IP tinggi, Minat rendah).
Meskipun IP di atas 3.5, area ini tetap MERAH/KUNING.
Kenapa? Karena Rule R7 membatasi output jadi Sedang untuk IP bagus × Minat rendah.
TANPA Rule R7, seluruh sisi kanan akan hijau — dan SKJK akan menang.
R7 membuat "jurang" di area kanan bawah, memastikan minat tetap diperhitungkan.

Ini membuktikan bahwa fuzzy system menangkap kedua faktor secara seimbang,
bukan hanya mengambil IP tertinggi.""")

# ═════════════════════════════════════════════════════════════════════
# SLIDE 10: Kesimpulan
# ═════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "Kesimpulan", font_size=30, bold=True, color=WHITE)
add_accent_line(slide, 0.8, 1.1, 3)

conclusions = [
    ("01", "Fuzzifikasi", "3 MF segitiga per variabel.\nSederhana, efisien, overlap halus.", BLUE),
    ("02", "9 Fuzzy Rules", "R7 = differentiator kunci:\nIP bagus tanpa minat → hanya Sedang.", GOLD),
    ("03", "Mamdani + COG", "Area fuzzy divisualisasikan.\nCOG = titik pusat massa area.", GREEN),
    ("04", "Verifikasi", "19/19 data latih = 100% akurat.\nTarget #20 → RPLD ✓", ACCENT),
]

for i, (num, title, desc, color) in enumerate(conclusions):
    y = 1.5 + i * 1.4
    add_card(slide, 0.8, y, 11.5, 1.2, BG_CARD)

    # Number
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), Inches(y + 0.2), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, 2.1, y + 0.1, 3, 0.5, title, font_size=18, bold=True, color=color)
    add_textbox(slide, 2.1, y + 0.55, 10, 0.6, desc, font_size=14, color=LIGHT)

# Bottom: Demo link + closing
add_card(slide, 1.0, 6.85, 11.3, 0.55, BG_CARD)
add_textbox(slide, 1.2, 6.9, 8, 0.4,
    "Demo interaktif: https://cpmk-02-ririx.streamlit.app/",
    font_size=13, color=GOLD)
add_textbox(slide, 9.5, 6.9, 2.7, 0.4,
    "Terima kasih!",
    font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

add_notes(slide, """HINT — SLIDE 10 (Kesimpulan) — ~1 menit

4 poin kesimpulan:

1. FUZZIFIKASI: 3 MF segitiga per variabel sudah cukup.
   Kenapa segitiga? Sederhana, komputasi ringan, overlap halus.

2. 9 FUZZY RULES: Yang paling penting adalah R7.
   R7 memastikan bahwa IP bagus saja TIDAK CUKUP tanpa minat.
   Ini yang membedakan fuzzy dari sekadar "ambil IP tertinggi".

3. MAMDANI + COG: Metode Mamdani menghasilkan area fuzzy
   yang bisa divisualisasikan, cocok untuk SPK yang perlu transparansi.
   COG = titik pusat massa dari area fuzzy.

4. VERIFIKASI: 100% akurat pada data latih, dan target #20
   berhasil merekomendasikan RPLD sesuai syarat soal ujian.

Sekian presentasi saya. Terima kasih. Ada pertanyaan?

---

PERSIAPAN ANTISIPASI PERTANYAAN DOSEN:

Q: Kenapa tidak pakai Sugeno?
A: Sugeno tidak menghasilkan area fuzzy, hanya angka. Untuk SPK
   yang perlu transparansi dan penjelasan ke pengguna, Mamdani
   lebih cocok karena setiap langkah bisa divisualisasikan.

Q: Kenapa R7 bukan Tinggi?
A: Pengalaman menunjukkan mahasiswa dengan IP tinggi tapi minat
   rendah cenderung drop-out atau tidak puas. Minat adalah faktor
   keberlanjutan, bukan hanya kemampuan awal.

Q: Kenapa parameter MF begitu?
A: Parameter disesuaikan dengan sebaran data survei. IP Bagus
   dimulai dari 3 karena itu batas "baik" di perguruan tinggi.
   Minat Suka berpuncak di 3 agar ada selisih jelas antara
   minat 2.60 dan 2.80 untuk data target #20.

Q: Bagaimana kalau datanya berubah?
A: Parameter MF dan rules bisa di-tuning menggunakan data training.
   Sistem ini bisa diupgrade dengan learning otomatis (ANFIS).""")

# ═════════════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════════════
output_path = os.path.join(OUTDIR, "presentasi_cpmk2.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
